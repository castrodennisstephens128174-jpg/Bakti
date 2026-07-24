#!/usr/bin/env python3
"""Build the nine-slide Bakti pitch deck from slides/slides.md."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
SOURCE = BASE / "slides.md"
OUTPUT = BASE / "Bakti.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BRAND = RGBColor(0x03, 0x69, 0xA1)
BRAND_SOFT = RGBColor(0xE6, 0xF3, 0xFB)
BRAND_LINE = RGBColor(0x9D, 0xCB, 0xEA)
INK = RGBColor(0x0B, 0x1B, 0x2B)
INK_SOFT = RGBColor(0x51, 0x61, 0x7A)
ACCENT = RGBColor(0xFD, 0xF1, 0xE0)
LINE = RGBColor(0xDB, 0xE3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Fraunces"
FONT_BODY = "Manrope"

HEADERS = ["TITLE", "PERSONA", "EVIDENCE", "CUSTOMER", "PRODUCT", "FLOW", "WHY", "MODEL", "STATUS"]


def parse_source(path: Path) -> list[dict]:
    slides: list[dict] = []
    current: dict | None = None
    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line or line.startswith("# Bakti Pitch") or line.startswith("# Exactly"):
            continue
        if line.startswith("# "):
            header = line[2:].strip()
            if header not in HEADERS:
                continue
            if current:
                slides.append(current)
            current = {"type": header.lower(), "lines": []}
        elif current:
            current["lines"].append(line)
    if current:
        slides.append(current)
    if len(slides) != 9:
        raise ValueError(f"Expected exactly 9 slides, found {len(slides)}")
    return slides


def add_blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), BRAND, BRAND)
    return slide


def add_rect(slide, left, top, width, height, fill, line=LINE, radius=False, dashed=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    if dashed:
        shape.line.dash_style = 4
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=18,
    color=INK,
    bold=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, left, top, width, height, size=16, bullet=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(7)
        p.line_spacing = 1.05
        p.text = ("• " if bullet else "") + line
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = INK_SOFT
    return box


def add_header(slide, kicker, title, num):
    add_text(slide, kicker.upper(), Inches(0.72), Inches(0.34), Inches(10.5), Inches(0.32), 11, BRAND, True)
    add_text(slide, title, Inches(0.72), Inches(0.73), Inches(11.8), Inches(0.72), 27, INK, True, FONT_HEAD)
    add_text(slide, f"{num} / 9", Inches(11.95), Inches(0.32), Inches(0.62), Inches(0.25), 9, INK_SOFT, False, FONT_BODY, PP_ALIGN.RIGHT)


def add_tag(slide, text, left, top, width, future=False):
    add_rect(slide, left, top, width, Inches(0.34), WHITE if future else BRAND_SOFT, BRAND_LINE, radius=True, dashed=future)
    add_text(slide, text, left, top + Inches(0.02), width, Inches(0.25), 9, BRAND, True, FONT_BODY, PP_ALIGN.CENTER)


def build_title(slide, data, num):
    lines = data["lines"]
    add_text(slide, "HOP STELLAR 2026 · PRODUCT PITCH", Inches(0.75), Inches(0.7), Inches(8), Inches(0.3), 12, BRAND, True)
    add_text(slide, lines[0], Inches(0.75), Inches(1.15), Inches(8.5), Inches(1.0), 54, BRAND, True, FONT_HEAD)
    add_text(slide, lines[1], Inches(0.75), Inches(2.15), Inches(11), Inches(0.65), 27, INK, True, FONT_HEAD)
    add_text(slide, lines[2], Inches(0.75), Inches(2.95), Inches(11), Inches(0.45), 18, INK_SOFT)
    add_rect(slide, Inches(0.75), Inches(3.8), Inches(5.75), Inches(1.65), BRAND_SOFT, BRAND_LINE, radius=True)
    add_tag(slide, "CURRENT", Inches(1.02), Inches(4.03), Inches(1.0))
    add_text(slide, "Stellar testnet prototype", Inches(1.02), Inches(4.5), Inches(5), Inches(0.32), 18, INK, True)
    add_text(slide, "Signed XLM escrow releases, direct payments, and on-chain verification.", Inches(1.02), Inches(4.9), Inches(5), Inches(0.42), 13, INK_SOFT)
    add_rect(slide, Inches(6.82), Inches(3.8), Inches(5.75), Inches(1.65), WHITE, BRAND_LINE, radius=True, dashed=True)
    add_tag(slide, "PLANNED", Inches(7.09), Inches(4.03), Inches(1.0), True)
    add_text(slide, "Licensed PHP cash-out", Inches(7.09), Inches(4.5), Inches(5), Inches(0.32), 18, INK, True)
    add_text(slide, "Provider SEP-24, KYC, routing/status, and confirmed collection.", Inches(7.09), Inches(4.9), Inches(5), Inches(0.42), 13, INK_SOFT)
    add_text(slide, lines[3], Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.25), 10, INK_SOFT)
    add_text(slide, f"{num} / 9", Inches(11.95), Inches(0.32), Inches(0.62), Inches(0.25), 9, INK_SOFT, False, FONT_BODY, PP_ALIGN.RIGHT)


def build_persona(slide, data, num):
    add_header(slide, data["lines"][0], "“I want support to be ready around salary day.”", num)
    add_tag(slide, data["lines"][1].upper(), Inches(0.75), Inches(1.62), Inches(3.2))
    body = data["lines"][2:]
    add_text(slide, body[0], Inches(0.75), Inches(2.2), Inches(7.2), Inches(1.75), 23, INK, False, FONT_HEAD)
    add_text(slide, body[1], Inches(0.75), Inches(4.2), Inches(7.2), Inches(0.78), 16, INK_SOFT)
    add_text(slide, body[2], Inches(0.75), Inches(5.25), Inches(7.2), Inches(0.55), 19, BRAND, True)
    add_rect(slide, Inches(8.65), Inches(1.6), Inches(3.9), Inches(4.9), ACCENT, LINE, radius=True)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.75), Inches(2.45), Inches(1.7), Inches(1.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = BRAND; circle.line.fill.background()
    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.05), Inches(2.72), Inches(1.1), Inches(0.85))
    face.fill.solid(); face.fill.fore_color.rgb = ACCENT; face.line.fill.background()
    add_text(slide, "MARIA", Inches(9.35), Inches(4.55), Inches(2.5), Inches(0.4), 12, INK_SOFT, True, FONT_BODY, PP_ALIGN.CENTER)
    add_text(slide, "Illustrative persona", Inches(9.35), Inches(5.02), Inches(2.5), Inches(0.35), 13, INK, True, FONT_BODY, PP_ALIGN.CENTER)
    add_text(slide, "Not a claimed interview", Inches(9.35), Inches(5.42), Inches(2.5), Inches(0.35), 11, INK_SOFT, False, FONT_BODY, PP_ALIGN.CENTER)


def build_evidence(slide, data, num):
    add_header(slide, data["lines"][0], "A large national flow, with a measured Malaysia signal", num)
    stats = data["lines"][1:4]
    for i, line in enumerate(stats):
        value, label = [part.strip() for part in line.split("|", 1)]
        left = Inches(0.75 + i * 4.08)
        add_rect(slide, left, Inches(1.7), Inches(3.75), Inches(1.75), ACCENT, ACCENT, radius=True)
        add_text(slide, value, left + Inches(0.22), Inches(1.98), Inches(3.3), Inches(0.55), 25, INK, True, FONT_HEAD)
        add_text(slide, label, left + Inches(0.22), Inches(2.62), Inches(3.3), Inches(0.6), 12, INK_SOFT)
    add_rect(slide, Inches(0.75), Inches(3.8), Inches(11.82), Inches(1.6), WHITE, LINE, radius=True)
    add_text(slide, data["lines"][4], Inches(1.0), Inches(4.12), Inches(11.2), Inches(0.5), 16, INK, True)
    add_text(slide, data["lines"][5], Inches(1.0), Inches(4.73), Inches(11.2), Inches(0.42), 13, INK_SOFT)
    add_text(slide, data["lines"][6], Inches(0.75), Inches(6.65), Inches(11.8), Inches(0.25), 8.5, INK_SOFT)


def build_customer(slide, data, num):
    add_header(slide, data["lines"][0], "One narrow customer; one clear job to be done", num)
    rows = data["lines"][1:]
    top = 1.62
    for idx, line in enumerate(rows):
        key, value = [part.strip() for part in line.split("|", 1)]
        fill = BRAND_SOFT if idx % 2 == 0 else WHITE
        add_rect(slide, Inches(0.75), Inches(top + idx * 0.67), Inches(11.82), Inches(0.58), fill, LINE, radius=True)
        add_text(slide, key, Inches(0.98), Inches(top + 0.13 + idx * 0.67), Inches(2.45), Inches(0.28), 12, INK, True)
        add_text(slide, value, Inches(3.22), Inches(top + 0.13 + idx * 0.67), Inches(9.0), Inches(0.3), 12, INK_SOFT)


def split_sections(lines: list[str]) -> tuple[list[str], list[str]]:
    marker = lines.index("NEXT — PLANNED")
    return lines[1:marker], lines[marker + 1:]


def build_product(slide, data, num):
    add_header(slide, data["lines"][0], "Do not confuse on-chain proof with cash delivery", num)
    current, future = split_sections(data["lines"][1:])
    add_rect(slide, Inches(0.75), Inches(1.65), Inches(5.72), Inches(4.75), BRAND_SOFT, BRAND_LINE, radius=True)
    add_tag(slide, "TODAY · WORKING", Inches(1.02), Inches(1.9), Inches(1.62))
    add_rich_lines(slide, current, Inches(1.02), Inches(2.48), Inches(5.05), Inches(3.6), 13.5)
    add_rect(slide, Inches(6.85), Inches(1.65), Inches(5.72), Inches(4.75), WHITE, BRAND_LINE, radius=True, dashed=True)
    add_tag(slide, "NEXT · PLANNED", Inches(7.12), Inches(1.9), Inches(1.6), True)
    add_rich_lines(slide, future, Inches(7.12), Inches(2.48), Inches(5.05), Inches(3.6), 13.5)


def build_flow(slide, data, num):
    add_header(slide, data["lines"][0], "Solid is current. Dashed is planned.", num)
    labels = [
        ("Sender", "Filipino worker\nin Malaysia"),
        ("Bakti plan", "Recipient + amount\n+ planning date"),
        ("Stellar", "Direct payment or\nXLM escrow release"),
        ("Recipient wallet", "Entered Stellar\naddress"),
    ]
    for idx, (title, subtitle) in enumerate(labels):
        left = 0.72 + idx * 3.14
        add_rect(slide, Inches(left), Inches(2.05), Inches(2.25), Inches(1.55), WHITE, BRAND_LINE, radius=True)
        add_text(slide, title, Inches(left + 0.12), Inches(2.37), Inches(2.0), Inches(0.3), 16, INK, True, FONT_BODY, PP_ALIGN.CENTER)
        add_text(slide, subtitle, Inches(left + 0.12), Inches(2.82), Inches(2.0), Inches(0.55), 11, INK_SOFT, False, FONT_BODY, PP_ALIGN.CENTER)
        if idx < 3:
            add_text(slide, "→", Inches(left + 2.31), Inches(2.53), Inches(0.75), Inches(0.4), 25, BRAND, True, FONT_BODY, PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.75), Inches(4.25), Inches(11.82), Inches(1.2), WHITE, BRAND_LINE, radius=True, dashed=True)
    add_tag(slide, "PLANNED", Inches(1.0), Inches(4.52), Inches(1.0), True)
    add_text(slide, data["lines"][2], Inches(2.2), Inches(4.48), Inches(9.95), Inches(0.48), 14, INK, True)
    add_text(slide, data["lines"][3], Inches(0.95), Inches(5.75), Inches(11.4), Inches(0.35), 12, INK_SOFT, False, FONT_BODY, PP_ALIGN.CENTER)
    add_text(slide, data["lines"][4], Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.35), 12, INK_SOFT, False, FONT_BODY, PP_ALIGN.CENTER)


def build_why(slide, data, num):
    add_header(slide, data["lines"][0], "Verifiable rails now; regulated last mile next", num)
    left_lines = data["lines"][1:6]
    right_lines = data["lines"][6:8]
    add_rect(slide, Inches(0.75), Inches(1.65), Inches(5.72), Inches(4.8), BRAND_SOFT, BRAND_LINE, radius=True)
    add_text(slide, "WHY STELLAR NOW", Inches(1.02), Inches(1.96), Inches(2.5), Inches(0.3), 12, BRAND, True)
    add_rich_lines(slide, left_lines, Inches(1.02), Inches(2.42), Inches(5.05), Inches(3.55), 13.5)
    add_rect(slide, Inches(6.85), Inches(1.65), Inches(5.72), Inches(4.8), WHITE, BRAND_LINE, radius=True, dashed=True)
    add_text(slide, "PROVIDER PATH", Inches(7.12), Inches(1.96), Inches(2.5), Inches(0.3), 12, BRAND, True)
    add_rich_lines(slide, right_lines, Inches(7.12), Inches(2.42), Inches(5.05), Inches(3.55), 12.5)
    add_text(slide, data["lines"][8], Inches(0.75), Inches(6.68), Inches(11.8), Inches(0.22), 8.5, INK_SOFT)


def build_model(slide, data, num):
    add_header(slide, data["lines"][0], "Hypotheses to test — not forecasts", num)
    lines = data["lines"][1:]
    marker = lines.index("GTM EXPERIMENTS")
    model = lines[1:marker]
    gtm = lines[marker + 1:]
    add_rect(slide, Inches(0.75), Inches(1.65), Inches(5.72), Inches(4.8), ACCENT, ACCENT, radius=True)
    add_tag(slide, "UNVALIDATED BUSINESS MODEL", Inches(1.02), Inches(1.92), Inches(2.38))
    add_rich_lines(slide, model, Inches(1.02), Inches(2.48), Inches(5.05), Inches(3.45), 14)
    add_rect(slide, Inches(6.85), Inches(1.65), Inches(5.72), Inches(4.8), WHITE, LINE, radius=True)
    add_tag(slide, "GTM EXPERIMENTS", Inches(7.12), Inches(1.92), Inches(1.72))
    add_rich_lines(slide, gtm, Inches(7.12), Inches(2.48), Inches(5.05), Inches(3.45), 13.5)


def build_status(slide, data, num):
    add_header(slide, data["lines"][0], "A working testnet core, with an honest last-mile gap", num)
    lines = data["lines"][1:]
    not_built = lines.index("NOT BUILT")
    ask = lines.index("ASK")
    built = lines[1:not_built]
    missing = lines[not_built + 1:ask]
    asks = lines[ask + 1:-1]
    proof = lines[-1]
    add_rect(slide, Inches(0.75), Inches(1.65), Inches(5.72), Inches(4.85), BRAND_SOFT, BRAND_LINE, radius=True)
    add_tag(slide, "BUILT", Inches(1.02), Inches(1.92), Inches(0.78))
    add_rich_lines(slide, built[:3], Inches(1.02), Inches(2.42), Inches(5.05), Inches(1.92), 12.2)
    add_text(slide, built[0].split(": ", 1)[1], Inches(1.02), Inches(4.45), Inches(5.05), Inches(0.55), 9.5, INK_SOFT, False, "Menlo")
    add_text(slide, built[1].split(": ", 1)[1], Inches(1.02), Inches(5.2), Inches(5.05), Inches(0.55), 9.5, INK_SOFT, False, "Menlo")
    add_rect(slide, Inches(6.85), Inches(1.65), Inches(5.72), Inches(4.85), WHITE, BRAND_LINE, radius=True, dashed=True)
    add_tag(slide, "ASK", Inches(7.12), Inches(1.92), Inches(0.72), True)
    add_rich_lines(slide, asks, Inches(7.12), Inches(2.42), Inches(5.05), Inches(2.05), 13.2)
    add_text(slide, "NOT BUILT", Inches(7.12), Inches(4.68), Inches(1.3), Inches(0.3), 11, BRAND, True)
    add_text(slide, missing[0], Inches(7.12), Inches(5.08), Inches(5.05), Inches(0.65), 11.5, INK_SOFT)
    add_text(slide, proof, Inches(0.75), Inches(6.72), Inches(11.8), Inches(0.25), 8.5, INK_SOFT)


def build_deck():
    data = parse_source(SOURCE)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    builders = {
        "title": build_title,
        "persona": build_persona,
        "evidence": build_evidence,
        "customer": build_customer,
        "product": build_product,
        "flow": build_flow,
        "why": build_why,
        "model": build_model,
        "status": build_status,
    }
    for num, item in enumerate(data, start=1):
        slide = add_blank(prs)
        builders[item["type"]](slide, item, num)
    prs.save(OUTPUT)
    print(f"Built {OUTPUT} — {len(prs.slides)} slides — {OUTPUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build_deck()
